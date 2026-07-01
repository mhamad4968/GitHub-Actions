# -*- coding: utf-8 -*-
"""2026年度 情報セキュリティ勉強会 PowerPoint（全12ページ構成）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import os
import sys

_REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
OUT = os.path.join(_REPO, "docs", "training", "security", "output", "情報セキュリティ勉強会テキスト_2026.pptx")

C_HEAD = RGBColor(0x1A, 0x36, 0x5D)
C_BODY = RGBColor(0x3A, 0x40, 0x45)
C_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
C_ACCENT = RGBColor(0xC0, 0x00, 0x00)
C_RED = RGBColor(0xE6, 0x30, 0x54)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN = RGBColor(0x2E, 0x7D, 0x32)
C_YELLOW = RGBColor(0xF9, 0xA8, 0x25)
C_ORANGE = RGBColor(0xE6, 0x51, 0x00)
C_LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
C_LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
C_LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
C_LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
C_WARN = RGBColor(0xFF, 0xF3, 0xCD)
C_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
FONT = "游ゴシック"
W = Inches(13.333)
H = Inches(7.5)


def run(r, t, sz=14, b=False, c=C_BODY):
    r.text = t
    r.font.name = FONT
    r.font.size = Pt(sz)
    r.font.bold = b
    r.font.color.rgb = c


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_tag(slide, n):
    box = slide.shapes.add_textbox(Inches(12.45), Inches(7.05), Inches(0.75), Inches(0.35))
    run(box.text_frame.paragraphs[0].add_run(), f"P{n:02d}", 11)


def title_bar(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.85))
    tf = box.text_frame
    run(tf.paragraphs[0].add_run(), text, 24, True, C_HEAD)
    if sub:
        p = tf.add_paragraph()
        run(p.add_run(), sub, 13, False, C_BLUE)


def body_box(slide, lines, left=0.55, top=1.0, width=12.2, height=5.8, sz=14):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        if isinstance(item, str):
            run(p.add_run(), item, sz)
        elif isinstance(item, tuple):
            run(p.add_run(), item[0], item[1] if len(item) > 1 else sz, item[2] if len(item) > 2 else False, item[3] if len(item) > 3 else C_BODY)
        elif isinstance(item, dict) and item.get("bullet"):
            run(p.add_run(), "● " + item["bullet"], sz)
        elif isinstance(item, list):
            for spec in item:
                if isinstance(spec, str):
                    run(p.add_run(), spec, sz)
                else:
                    run(p.add_run(), spec["t"], spec.get("sz", sz), spec.get("b", False), spec.get("c", C_BODY))


def rounded_block(slide, left, top, width, height, fill, line=None, text_lines=None, title=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(2)
    else:
        sh.line.fill.background()
    if text_lines or title:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_top = Pt(6)
        if title:
            run(tf.paragraphs[0].add_run(), title, 13, True, C_HEAD)
        for line in (text_lines or []):
            p = tf.add_paragraph()
            run(p.add_run(), line, 12)
    return sh


def flow_box(slide, left, top, w, h, text, fill=C_LIGHT_BLUE, accent=C_BLUE):
    sh = rounded_block(slide, left, top, w, h, fill, accent, [text])
    return sh


def arrow_text(slide, left, top, text="➡"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(0.5), Inches(0.4))
    run(box.text_frame.paragraphs[0].add_run(), text, 22, True, C_ACCENT)


def add_table_simple(slide, headers, rows, left, top, col_w):
    nr = 1 + len(rows)
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(sum(col_w)), Inches(0.42 * nr))
    tbl = ts.table
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C_HEAD
        run(c.text_frame.paragraphs[0].add_run(), h, 11, True, C_WHITE)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            run(tbl.cell(ri + 1, ci).text_frame.paragraphs[0].add_run(), val, 10, ci == 0, C_BODY)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ═══════════════════════════════════════
# P01 タイトル
# ═══════════════════════════════════════
s = blank(prs)
# 背景デコ：盾・鍵アイコン風
for x, y, sz in [(10.5, 0.8, 2.2), (1.0, 5.5, 1.8), (11.2, 4.8, 1.4)]:
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(sz), Inches(sz))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_LIGHT_BLUE
    sh.line.color.rgb = C_BLUE
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    icon = "🛡" if x > 5 else "🔒"
    run(tf.paragraphs[0].add_run(), icon, 36 if sz > 2 else 28)

box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.2))
tf = box.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0].add_run(), "2026年度 情報セキュリティ勉強会", 38, True, C_HEAD)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run(p2.add_run(), "テキスト", 30, True, C_HEAD)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run(p3.add_run(), "最新のサイバー脅威と私たちが果たすべき防衛策", 18, False, C_BLUE)
p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(20)
run(p4.add_run(), "システム推進室 / 2026.7.15", 14, False, C_BODY)

corner = s.shapes.add_textbox(Inches(9.8), Inches(6.55), Inches(3.2), Inches(0.5))
run(corner.text_frame.paragraphs[0].add_run(), "システム推進室", 12, False, C_BODY)
page_tag(s, 1)

# ═══════════════════════════════════════
# P02 アジェンダ（垂直ブロックリスト風）
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "本日のお題")
agenda = [
    ("1", "情報セキュリティをめぐる最新動向", C_HEAD),
    ("2", "組織を狙う最凶の罠「ランサムウェア」", C_ACCENT),
    ("3", "身近に潜む罠「偽セキュリティ（サポート詐欺）」", C_ORANGE),
    ("4", "生成AIの悪用「超巧妙な標的型メール」", C_BLUE),
    ("5", "被害に遭わないための4つの鉄則", C_GREEN),
]
for i, (num, txt, col) in enumerate(agenda):
    top = 1.15 + i * 1.08
    rounded_block(s, 0.7, top, 11.9, 0.92, C_WHITE, col, [f"{num}.  {txt}"], None)
    # 左色帯
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(top), Inches(0.18), Inches(0.92))
    band.fill.solid()
    band.fill.fore_color.rgb = col
    band.line.fill.background()
body_box(s, [("研修時間目安：テキスト 20分 ／ 動画視聴 20分", 13, True, C_BODY)], top=6.55, height=0.4)
page_tag(s, 2)

# ═══════════════════════════════════════
# P03 最新動向 + IPA表
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "1. 情報セキュリティをめぐる最新動向")
body_box(s, [
    "サイバー攻撃は常に「いたちごっこ」",
    "システム対策の強化だけでは防ぎきれません。",
    [{"t": "「社員一人ひとりの意識」", "b": True, "c": C_ACCENT, "sz": 15},
     {"t": "こそが、組織を守る最後の砦です。", "sz": 15}],
], left=0.55, top=1.05, width=5.8, height=2.5, sz=14)

rounded_block(s, 0.55, 3.7, 5.8, 1.2, C_LIGHT_RED, C_ACCENT, [
    "2020年 日経グループ：1台の感染から約1.2万人分の個人情報流出",
    "検知システムがあっても、新種の手口では検知が遅れる場合あり",
])

title_bar(s, "", "最新の脅威トレンド（IPA調べ）")
add_table_simple(s, ["順位", "組織向け脅威", "備考"], [
    ["1位", "ランサム攻撃による被害", "11年連続"],
    ["2位", "サプライチェーン・委託先を狙った攻撃", "8年連続"],
    ["3位", "AIの利用をめぐるサイバーリスク", "2026年初選出"],
    ["4位", "システムの脆弱性を悪用した攻撃", "6年連続"],
    ["5位", "機密情報等を狙った標的型攻撃", "11年連続"],
], 6.7, 1.05, [0.7, 3.8, 1.5])

body_box(s, [
    [{"t": "🔗 IPA公式：", "sz": 10}, {"t": "情報セキュリティ10大脅威 2026", "sz": 10, "c": C_BLUE}],
    [{"t": "https://www.ipa.go.jp/security/10threats/10threats2026.html", "sz": 9, "c": C_BLUE}],
], left=6.7, top=5.5, width=5.8, height=0.8, sz=10)
page_tag(s, 3)

# ═══════════════════════════════════════
# P04 ランサムウェア + マルチ脅迫フロー
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "2. 組織を狙う最凶の罠「ランサムウェア」")
body_box(s, [
    [{"t": "ランサムウェア（身代金要求型不正プログラム）とは？", "b": True, "c": C_HEAD, "sz": 15}],
    {"bullet": "PC内のファイルやシステムを勝手に暗号化して使えなくする。"},
    {"bullet": "元に戻すことと引き換えに金銭（身代金）を要求。"},
    [{"t": "「暗号化」だけではない！進化するマルチ脅迫", "b": True, "c": C_ACCENT, "sz": 14}],
    {"bullet": "【脅迫1】データを暗号化してシステムを人質にする。"},
    {"bullet": "【脅迫2】事前に盗み出した機密データをネット上に暴露すると脅す。"},
], left=0.55, top=1.0, width=12.0, height=2.8, sz=13)

# フロー図
flow_box(s, 0.8, 4.2, 3.2, 0.85, "① 感染\n（メール・脆弱性等）", C_LIGHT_RED, C_ACCENT)
arrow_text(s, 4.15, 4.45)
flow_box(s, 4.6, 4.2, 3.2, 0.85, "② ファイル暗号化\nシステム使用不能", C_LIGHT_YELLOW, C_ORANGE)
arrow_text(s, 7.95, 4.45)
flow_box(s, 8.4, 4.2, 3.5, 0.85, "③ 身代金要求\n＋ データ暴露脅迫", C_LIGHT_RED, C_ACCENT)

body_box(s, [
    "※ 三重・四重脅迫：DDoS攻撃、顧客・取引先への連絡等も報告されています。",
], left=0.55, top=5.35, width=12.0, height=0.5, sz=11)
page_tag(s, 4)

# ═══════════════════════════════════════
# P05 感染経路 円グラフ
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "【データ】どこから感染する？ランサムウェアの経路")
chart_data = CategoryChartData()
chart_data.categories = [
    "通信機器の設定不備・脆弱性",
    "リモートデスクトップ侵入",
    "不審メール・添付",
    "その他",
]
chart_data.add_series("感染経路", (62, 19, 9, 10))
chart_frame = s.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(0.4), Inches(1.1), Inches(5.8), Inches(5.2), chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.plots[0].has_data_labels = True

body_box(s, [
    [{"t": "日本国内のランサムウェア感染経路（警察庁調べ）", "b": True, "sz": 14, "c": C_HEAD}],
    ("🔴 通信機器（VPNなど）の設定不備・脆弱性：62%", 13),
    ("🔵 リモートデスクトップからの侵入：19%", 13),
    ("🟢 不審なメールやその添付ファイル：9%", 13),
    ("⚪ その他：10%", 13),
], left=6.5, top=1.2, width=6.3, height=2.5, sz=13)

rounded_block(s, 6.5, 4.0, 6.3, 2.0, C_WARN, C_ORANGE, None)
warn_tf = s.shapes[-1].text_frame
run(warn_tf.paragraphs[0].add_run(), "💡 重要な事実", 14, True, C_ACCENT)
p = warn_tf.add_paragraph()
run(p.add_run(), "「怪しいメールを開かなければ安全」という時代は終わりました。", 12, True, C_BODY)
p2 = warn_tf.add_paragraph()
run(p2.add_run(), "設定の甘い機器からネットワークの隙を突いて裏から侵入されるケースが圧倒的多数です。", 12, False, C_BODY)
page_tag(s, 5)

# ═══════════════════════════════════════
# P06 KADOKAWA 事例①
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "【事例①】KADOKAWAランサムウェア事件の概要")
body_box(s, [
    [{"t": "事件の概要（2024年6月）", "b": True, "c": C_HEAD, "sz": 15}],
    {"bullet": "攻撃者：ロシア系サイバー犯罪集団「BlackSuit」"},
    {"bullet": "手法：ランサムウェア攻撃 ＋ 約1.5TBのデータ窃取"},
    [{"t": "恐るべき「侵入経路」", "b": True, "c": C_ACCENT, "sz": 14}],
], left=0.55, top=1.0, width=5.8, height=2.5, sz=13)

steps = [
    "① フィッシング詐欺等で、\nたった一人の従業員の\nアカウント情報が盗まれる",
    "② 攻撃者が「正規の社員」の\nふりをして社内ネットワーク\nへ堂々と侵入",
    "③ 内部で権限を拡大し、\nデータセンターのサーバー群\nを一斉暗号化",
]
for i, st in enumerate(steps):
    flow_box(s, 6.6, 1.1 + i * 1.85, 5.9, 1.55, st, C_LIGHT_BLUE, C_BLUE)
    if i < 2:
        arrow_text(s, 9.3, 2.75 + i * 1.85, "▼")

rounded_block(s, 0.55, 5.5, 12.0, 1.1, C_LIGHT_RED, C_ACCENT, [
    "教訓：強固なシステムがあっても、「たった一人のID流出」で組織全体が崩壊する。",
])
page_tag(s, 6)

# ═══════════════════════════════════════
# P07 KADOKAWA 事例②
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "【事例②】KADOKAWA事件の被害とデータ管理の落とし穴")
body_box(s, [
    [{"t": "甚大なビジネス被害", "b": True, "c": C_HEAD, "sz": 15}],
    {"bullet": "ニコニコ動画などのサービスが約2ヶ月間にわたり完全停止"},
    {"bullet": "出版物の出荷システムがダウンし、出荷数が平常時の3分の1に激減"},
    {"bullet": "25万人分以上の個人情報・機密契約書が流出"},
], left=0.55, top=1.0, width=12.0, height=2.3, sz=13)

warn = rounded_block(s, 0.55, 3.5, 12.0, 3.2, C_WARN, C_ACCENT, None)
tf = s.shapes[-1].text_frame
run(tf.paragraphs[0].add_run(), "⚠️ 安全なはずのデータがなぜ漏洩した？「放置コピーの罠」", 15, True, C_ACCENT)
for line in [
    "当初安全とされていた生徒情報（N高生等）が流出。",
    "原因：「特定の社員1名が、業務のために個人情報を別フォルダにコピーし、そのまま削除せず放置していた」",
    "教訓：利便性のためにデータを勝手にコピーし、放置する行為は致命的なリスクになる。",
]:
    p = tf.add_paragraph()
    run(p.add_run(), line, 13, "教訓" in line, C_BODY if "教訓" not in line else C_ACCENT)
page_tag(s, 7)

# ═══════════════════════════════════════
# P08 偽セキュリティ
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "3. 身近に潜む罠「偽セキュリティ（サポート詐欺）」")
steps8 = [
    ("ステップ1", "偽の警告表示", "ネット閲覧中に突然、大音量の警告音と共に「ウイルス検出！」の画面を出す。"),
    ("ステップ2", "偽サポートへ誘導", "画面に「マイクロソフト」等のロゴを悪用し、偽の電話番号へ連絡させる。"),
    ("ステップ3", "遠隔操作・金銭要求", "電話をかけると、遠隔操作ソフトを入れさせられ、電子マネー等でサポート代金を騙し取る。"),
]
for i, (st, tit, desc) in enumerate(steps8):
    top = 1.1 + i * 1.65
    rounded_block(s, 0.55, top, 6.0, 1.45, C_WHITE, C_ORANGE, [desc], f"{st}：{tit}")

# 偽警告画面プレースホルダ
fake = rounded_block(s, 6.9, 1.1, 5.6, 4.8, RGBColor(0x1E, 0x1E, 0x1E), C_RED, None)
tf = s.shapes[-1].text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0].add_run(), "⚠ ウイルス検出！", 22, True, C_RED)
for t in ["お使いのPCは感染しています", "今すぐお電話ください", "0120-XXX-XXX  ✕", "（正規サポートは電話を促しません）"]:
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    c = C_RED if "✕" in t else C_WHITE
    run(p.add_run(), t, 14 if "0120" in t else 12, "0120" in t, c)

body_box(s, [
    [{"t": "📷 参考：警察庁 サポート詐欺対策特設ページ（偽画面イメージ）", "sz": 10, "c": C_BLUE}],
    [{"t": "※ 正規のMicrosoft・ウイルス対策は画面に電話番号を表示しません", "sz": 11, "b": True, "c": C_ACCENT}],
], left=0.55, top=6.15, width=12.0, height=0.7, sz=10)
page_tag(s, 8)

# ═══════════════════════════════════════
# P09 AI標的型メール 対比表
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "4. 生成AIの悪用「超巧妙な標的型メール」")
body_box(s, [
    "攻撃者は「生成AI（ChatGPTなど）」を悪用しています。",
    [{"t": "「目視の違和感だけでは、もう100%見抜けない」時代です。", "b": True, "c": C_ACCENT, "sz": 15}],
], left=0.55, top=0.95, width=12.0, height=1.0, sz=14)

add_table_simple(s, ["チェックポイント", "過去の不審メール ❌", "現在のAI悪用メール ⚠️"], [
    ["日本語の文体", "カタコト、文法が不自然", "完璧で極めて自然なビジネス敬語"],
    ["メールの文脈", "突然、見覚えのない内容が届く", "過去の本物のやり取りを学習し、業務の文脈を再現"],
    ["なりすまし対象", "海外の不審なアカウント", "本物の社長、上司、取引先の口調を完コピ"],
], 0.55, 2.15, [2.5, 4.5, 5.0])

rounded_block(s, 0.55, 5.0, 12.0, 1.5, C_LIGHT_BLUE, C_BLUE, [
    "対策：「いつもと違う」「急かす」「金銭・認証情報」のキーワードで止まる",
    "不審メールは削除せず、システム推進室へ報告（社内周知に役立つ）",
])
page_tag(s, 9)

# ═══════════════════════════════════════
# P10 4つの鉄則
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "まとめ：被害に遭わないための「4つの鉄則」")
rules = [
    ("🟩 鉄則1：開かない", "違和感のあるリンク・添付は即座に開かない\n少しでも怪しければ、開く前に電話など別ルートで相手に確認。", C_LIGHT_GREEN, C_GREEN),
    ("🟨 鉄則2：すぐ遮断", "怪しい画面・挙動時は「まずネットを切る」\nLANケーブルを抜く、Wi-Fiを切る。他PCへの感染拡大を防ぐのが最優先。", C_LIGHT_YELLOW, C_ORANGE),
    ("🟥 鉄則3：ルール厳守", "私物メディア・フリーWi-Fiの利用禁止\n未許可のUSBメモリは使わない。外出先の公共Wi-Fiには絶対に繋がない。", C_LIGHT_RED, C_ACCENT),
    ("🟦 鉄則4：即報告", "一人で悩まず「システム推進室」へ連絡\n「操作ミスで怪しいボタンを押してしまった」時こそ、隠さず即報告。", C_LIGHT_BLUE, C_BLUE),
]
for i, (tit, desc, fill, line) in enumerate(rules):
    col = i % 2
    row = i // 2
    rounded_block(s, 0.55 + col * 6.15, 1.15 + row * 2.85, 5.95, 2.55, fill, line, desc.split("\n"), tit)
page_tag(s, 10)

# ═══════════════════════════════════════
# P11 初動対応マニュアル
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "【初動対応マニュアル】警告が出たら、まずこれ！")

# 手順1
num1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(1.15), Inches(0.7), Inches(0.7))
num1.fill.solid()
num1.fill.fore_color.rgb = C_ACCENT
num1.line.fill.background()
run(num1.text_frame.paragraphs[0].add_run(), "1", 24, True, C_WHITE)
num1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

rounded_block(s, 1.4, 1.1, 5.5, 2.5, C_LIGHT_YELLOW, C_ORANGE, [
    "有線LANの場合：ツメを押さえながらLANケーブルを抜く",
    "Wi-Fiの場合：画面右下のアイコンからWi-Fiを切断",
    "（またはPCを「機内モード」にする）",
], "手順1：ネットワークを物理的に遮断する")

# 手順1 イメージ枠
img1 = rounded_block(s, 1.4, 3.75, 5.5, 1.5, C_GRAY, C_BODY, ["【操作イメージ】LANケーブル抜去 / Wi-Fi OFF"], None)

# 手順2
num2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(1.15), Inches(0.7), Inches(0.7))
num2.fill.solid()
num2.fill.fore_color.rgb = C_BLUE
num2.line.fill.background()
run(num2.text_frame.paragraphs[0].add_run(), "2", 24, True, C_WHITE)
num2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

rounded_block(s, 7.85, 1.1, 5.0, 2.5, C_LIGHT_BLUE, C_BLUE, [
    "一人で解決しようと画面を操作すると、被害が拡大します。",
    "上司・先輩社員に声をかけてください。",
    "システム推進室へ連絡（社内掲示・マニュアル参照）",
], "手順2：システム推進室へ連絡する")

rounded_block(s, 7.85, 3.75, 5.0, 1.5, C_LIGHT_RED, C_ACCENT, [
    "指示を受けるまで、該当PCで業務を再開しない",
    "※ 画面の閉じ方は動画「偽セキュリティ警告画面に注意！」で確認",
])

body_box(s, [
    [{"t": "※ 個人の携帯番号ではなく、必ずシステム推進室の公式連絡先を使用してください。", "sz": 11, "c": C_ACCENT}],
], left=0.55, top=5.6, width=12.0, height=0.4, sz=11)
page_tag(s, 11)

# ═══════════════════════════════════════
# P12 動画視聴・クロージング
# ═══════════════════════════════════════
s = blank(prs)
title_bar(s, "続いて：動画視聴（20分）")
body_box(s, [
    "テキスト研修の後、以下3本を視聴してください。",
    [{"t": "① そこにある脅威〜組織を狙うランサムウェア攻撃〜", "b": True, "sz": 14}],
    [{"t": "https://youtu.be/TWqJ5P8oaUM?t=8", "sz": 11, "c": C_BLUE}],
    [{"t": "② 偽セキュリティ警告画面に注意！【画面削除方法】", "b": True, "sz": 14}],
    [{"t": "https://youtu.be/BV6GebTEiQI", "sz": 11, "c": C_BLUE}],
    [{"t": "③ 見えざるサイバー攻撃", "b": True, "sz": 14}],
    [{"t": "https://www.youtube.com/watch?v=ZZJ7VMJ5Btw", "sz": 11, "c": C_BLUE}],
], left=0.55, top=1.2, width=12.0, height=4.0, sz=13)

box = s.shapes.add_textbox(Inches(0), Inches(5.8), W, Inches(1.2))
tf = box.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run(tf.paragraphs[0].add_run(), "お疲れ様でした！", 28, True, C_RED)
p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run(p.add_run(), "システム推進室", 16, False, C_BODY)
page_tag(s, 12)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
try:
    prs.save(OUT)
    print("Saved:", OUT, "slides:", len(prs.slides))
except PermissionError:
    alt = OUT.replace(".pptx", "_12slides.pptx")
    prs.save(alt)
    print("Saved (alt):", alt, "slides:", len(prs.slides))
